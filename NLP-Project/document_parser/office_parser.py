# office_parser.py
try:
    import docx
except ImportError:
    import python

    -docx as docx

from openpyxl import load_workbook
from pptx import Presentation
import os
from base_models import ParsedDocument


class OfficeParser:
    def parse_word(self, file_path: str, doc_id: str) -> ParsedDocument:
        """解析Word文档"""
        try:
            doc = docx.Document(file_path)
            content = []
            tables = []
            images = []

            # 提取文本内容
            for para in doc.paragraphs:
                if para.text.strip():
                    content.append(para.text)

            # 提取表格
            for i, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_data:  # 只保留非空行
                        table_data.append(row_data)

                if table_data:
                    tables.append({
                        "table_index": i,
                        "header": table_data[0] if table_data else [],
                        "rows": table_data[1:] if len(table_data) > 1 else []
                    })

            # 提取文档结构
            structure = self._extract_word_structure(doc)

            # 提取元数据
            metadata = {
                "filename": os.path.basename(file_path),
                "file_size": os.path.getsize(file_path),
                "paragraph_count": len(content),
                "table_count": len(tables)
            }

            return ParsedDocument(
                doc_id=doc_id,
                doc_type="word",
                source_path=file_path,
                content="\n".join(content),
                metadata=metadata,
                tables=tables,
                images=images,
                structure=structure
            )

        except Exception as e:
            print(f"Word解析错误 {file_path}: {e}")
            return self._create_error_document(doc_id, file_path, "word")

    def parse_excel(self, file_path: str, doc_id: str) -> ParsedDocument:
        """解析Excel文档"""
        try:
            wb = load_workbook(file_path)
            content_parts = []
            tables = []

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                content_parts.append(f"=== 工作表: {sheet_name} ===")

                # 提取表格数据
                data = []
                for row in ws.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        cleaned_row = [str(cell) if cell is not None else "" for cell in row]
                        data.append(cleaned_row)

                if data:
                    tables.append({
                        "sheet_name": sheet_name,
                        "header": data[0],
                        "rows": data[1:]
                    })

                    # 生成文字描述
                    if len(data) > 1:
                        desc = f"工作表'{sheet_name}'包含{len(data) - 1}行数据，列标题: {', '.join(data[0])}"
                        content_parts.append(desc)

            metadata = {
                "filename": os.path.basename(file_path),
                "file_size": os.path.getsize(file_path),
                "sheet_count": len(wb.sheetnames),
                "table_count": len(tables)
            }

            return ParsedDocument(
                doc_id=doc_id,
                doc_type="excel",
                source_path=file_path,
                content="\n".join(content_parts),
                metadata=metadata,
                tables=tables,
                structure={"sheets": wb.sheetnames}
            )

        except Exception as e:
            print(f"Excel解析错误 {file_path}: {e}")
            return self._create_error_document(doc_id, file_path, "excel")

    def parse_ppt(self, file_path: str, doc_id: str) -> ParsedDocument:
        """解析PPT文档"""
        try:
            prs = Presentation(file_path)
            content_parts = []
            images = []

            for i, slide in enumerate(prs.slides):
                slide_content = []

                # 提取文本内容
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text)

                    # 处理图表描述
                    if hasattr(shape, "chart"):
                        chart_desc = f"图表: {shape.name}" if shape.name else "图表"
                        slide_content.append(chart_desc)
                        images.append({
                            "type": "chart",
                            "description": chart_desc,
                            "slide": i + 1
                        })

                if slide_content:
                    content_parts.append(f"=== 幻灯片 {i + 1} ===")
                    content_parts.extend(slide_content)

            metadata = {
                "filename": os.path.basename(file_path),
                "file_size": os.path.getsize(file_path),
                "slide_count": len(prs.slides)
            }

            return ParsedDocument(
                doc_id=doc_id,
                doc_type="ppt",
                source_path=file_path,
                content="\n".join(content_parts),
                metadata=metadata,
                images=images,
                structure={"slides": len(prs.slides)}
            )

        except Exception as e:
            print(f"PPT解析错误 {file_path}: {e}")
            return self._create_error_document(doc_id, file_path, "ppt")

    def _extract_word_structure(self, doc) -> Dict:
        """提取Word文档结构"""
        sections = []
        for i, para in enumerate(doc.paragraphs):
            if para.style.name.startswith('Heading'):
                sections.append({
                    "title": para.text,
                    "style": para.style.name,
                    "level": int(para.style.name.replace('Heading', '')) if para.style.name.replace('Heading',
                                                                                                    '').isdigit() else 1,
                    "position": i
                })
        return {"headings": sections}

    def _create_error_document(self, doc_id: str, file_path: str, doc_type: str) -> ParsedDocument:
        """创建错误文档"""
        return ParsedDocument(
            doc_id=doc_id,
            doc_type=doc_type,
            source_path=file_path,
            content=f"文档解析失败: {file_path}",
            metadata={
                "filename": os.path.basename(file_path),
                "error": True
            }
        )